from __future__ import annotations

import os
from pathlib import Path
from typing import Any

DOCS_PATH = Path(__file__).parent.parent / "data" / "documents"
DOCS_PATH.mkdir(parents=True, exist_ok=True)


def create_word(filename: str, content: str) -> str:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    doc = Document()

    title = doc.add_heading(filename, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.runs[0]
    run.font.color.rgb = RGBColor(0, 95, 142)

    doc.add_paragraph("─" * 50)

    for line in content.split("\n"):
        if line.strip():
            p = doc.add_paragraph(line)
            for run in p.runs:
                run.font.size = Pt(12)

    path = DOCS_PATH / f"{filename}.docx"
    doc.save(str(path))
    os.startfile(path)
    return f"Word creado: {filename}.docx"


def create_excel(filename: str, content: str) -> str:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = filename

    header_font = Font(color="FFFFFF", bold=True, size=12)
    header_fill = PatternFill(start_color="005F8E", end_color="005F8E", fill_type="solid")

    for i, line in enumerate(content.split("\n"), 1):
        cols = [c.strip() for c in line.split(",")]
        for j, val in enumerate(cols, 1):
            cell = ws.cell(row=i, column=j, value=val)
            if i == 1:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")

    for col in ws.columns:
        max_len = max((len(str(cell.value or "")) for cell in col), default=0)
        ws.column_dimensions[col[0].column_letter].width = max(max_len + 4, 12)

    path = DOCS_PATH / f"{filename}.xlsx"
    wb.save(path)
    os.startfile(path)
    return f"Excel creado: {filename}.xlsx"


def create_pdf(filename: str, content: str) -> str:
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    path = DOCS_PATH / f"{filename}.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=letter, leftMargin=72, rightMargin=72, topMargin=72, bottomMargin=72)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("CustomTitle", parent=styles["Title"], textColor=HexColor("#005F8E"), fontSize=20, spaceAfter=20)
    body_style = ParagraphStyle("CustomBody", parent=styles["Normal"], fontSize=12, leading=18)

    story = [Paragraph(filename, title_style), Spacer(1, 12)]
    for line in content.split("\n"):
        if line.strip():
            story.append(Paragraph(line, body_style))
            story.append(Spacer(1, 6))

    doc.build(story)
    os.startfile(path)
    return f"PDF creado: {filename}.pdf"


def create_pptx(filename: str, content: str) -> str:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # Paleta de colores profesional
    palettes = [
        {  # Azul oscuro / cyan — estilo TUTUS
            "bg": RGBColor(2, 8, 20),  # type: ignore[no-untyped-call]
            "accent": RGBColor(0, 200, 255),  # type: ignore[no-untyped-call]
            "title": RGBColor(255, 255, 255),  # type: ignore[no-untyped-call]
            "body": RGBColor(180, 230, 255),  # type: ignore[no-untyped-call]
            "bar": RGBColor(0, 150, 200),  # type: ignore[no-untyped-call]
        },
        {  # Oscuro profesional
            "bg": RGBColor(18, 18, 30),  # type: ignore[no-untyped-call]
            "accent": RGBColor(120, 80, 255),  # type: ignore[no-untyped-call]
            "title": RGBColor(255, 255, 255),  # type: ignore[no-untyped-call]
            "body": RGBColor(200, 190, 255),  # type: ignore[no-untyped-call]
            "bar": RGBColor(80, 40, 200),  # type: ignore[no-untyped-call]
        },
        {  # Verde tech
            "bg": RGBColor(5, 20, 10),  # type: ignore[no-untyped-call]
            "accent": RGBColor(0, 255, 120),  # type: ignore[no-untyped-call]
            "title": RGBColor(255, 255, 255),  # type: ignore[no-untyped-call]
            "body": RGBColor(180, 255, 210),  # type: ignore[no-untyped-call]
            "bar": RGBColor(0, 180, 80),  # type: ignore[no-untyped-call]
        },
    ]

    palette = palettes[0]  # TUTUS style por default

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    lines = [line.strip() for line in content.split("\n") if line.strip()]

    def set_bg(slide: Any, color: Any) -> None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_bar(slide: Any, palette: dict[str, Any]) -> None:
        """Barra decorativa superior."""
        from pptx.util import Inches

        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0),
            Inches(0),
            prs.slide_width,
            Inches(0.08),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = palette["accent"]
        bar.line.fill.background()

    def add_bottom_bar(slide: Any, palette: dict[str, Any]) -> None:
        bar = slide.shapes.add_shape(1, Inches(0), (prs.slide_height or Inches(7.5)) - Inches(0.06), prs.slide_width, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = palette["bar"]
        bar.line.fill.background()

    def add_text_box(
        slide: Any,
        text: str,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        font_size: int = 24,
        color: Any = None,
        bold: bool = False,
        align: Any = PP_ALIGN.LEFT,
    ) -> Any:
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return tx_box

    # ── SLIDE 1: TÍTULO ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide1, palette["bg"])
    add_bar(slide1, palette)
    add_bottom_bar(slide1, palette)

    # Línea vertical decorativa
    line = slide1.shapes.add_shape(1, Inches(1.2), Inches(1.5), Inches(0.05), Inches(4))
    line.fill.solid()
    line.fill.fore_color.rgb = palette["accent"]
    line.line.fill.background()

    # Título principal
    add_text_box(slide1, filename, Inches(1.6), Inches(1.8), Inches(10), Inches(1.5), font_size=48, color=palette["title"], bold=True)

    # Subtítulo
    if lines:
        first = lines[0].split("|")
        subtitle = first[1] if len(first) > 1 else first[0]
        add_text_box(slide1, subtitle, Inches(1.6), Inches(3.5), Inches(10), Inches(1), font_size=20, color=palette["accent"])

    # Badge TUTUS
    add_text_box(
        slide1,
        "◈ Generado por TUTUS",
        Inches(9),
        Inches(6.8),
        Inches(4),
        Inches(0.5),
        font_size=10,
        color=palette["bar"],
        align=PP_ALIGN.RIGHT,
    )

    # ── SLIDES DE CONTENIDO ──
    for i, line in enumerate(lines[1:], 1):
        parts = line.split("|")
        titulo = parts[0].strip()
        contenido = parts[1].strip() if len(parts) > 1 else ""

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, palette["bg"])
        add_bar(slide, palette)
        add_bottom_bar(slide, palette)

        # Número de slide
        add_text_box(slide, f"{i:02d}", Inches(0.3), Inches(0.8), Inches(1), Inches(1), font_size=36, color=palette["accent"], bold=True)

        # Línea separadora título
        sep = slide.shapes.add_shape(1, Inches(1.2), Inches(1.6), Inches(11.5), Inches(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = palette["accent"]
        sep.line.fill.background()

        # Título slide
        add_text_box(slide, titulo, Inches(1.2), Inches(0.7), Inches(11), Inches(1), font_size=30, color=palette["title"], bold=True)

        # Contenido — procesar bullets
        if contenido:
            bullet_items = contenido.split("·") if "·" in contenido else [contenido]
            y_pos: Any = Inches(2.0)

            for item in bullet_items:
                item = item.strip()
                if not item:
                    continue

                # Bullet point decorativo
                dot = slide.shapes.add_shape(1, Inches(1.2), y_pos + Inches(0.12), Inches(0.08), Inches(0.08))
                dot.fill.solid()
                dot.fill.fore_color.rgb = palette["accent"]
                dot.line.fill.background()

                add_text_box(slide, item, Inches(1.5), y_pos, Inches(11), Inches(0.6), font_size=16, color=palette["body"])
                y_pos += Inches(0.75)

    path = DOCS_PATH / f"{filename}.pptx"
    prs.save(str(path))
    os.startfile(path)
    return f"Presentación creada: {filename}.pptx"
