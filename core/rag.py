from __future__ import annotations

import hashlib
import logging
from pathlib import Path
from typing import Any

import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer

from core.config import CHUNK_OVERLAP, CHUNK_SIZE, DOCUMENTS_DIR, EMBEDDING_MODEL, RAG_DB_DIR, RAG_TOP_K

log = logging.getLogger("tutus.rag")


_embedder: SentenceTransformer | None = None
_chroma_client: Any = None
_collection: Any = None


def _get_embedder() -> SentenceTransformer:
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(EMBEDDING_MODEL)
    return _embedder


def _get_collection() -> Any:
    global _chroma_client, _collection
    if _chroma_client is None:
        RAG_DB_DIR.mkdir(parents=True, exist_ok=True)
        _chroma_client = chromadb.PersistentClient(
            path=str(RAG_DB_DIR),
            settings=Settings(anonymized_telemetry=False),
        )
    if _collection is None:
        _collection = _chroma_client.get_or_create_collection(
            name="tutus_docs",
            metadata={"hnsw:space": "cosine"},
        )
    return _collection


def _chunk_text(text: str, source: str) -> list[dict[str, Any]]:
    words = text.split()
    chunks: list[dict[str, Any]] = []
    if not words:
        return chunks

    start = 0
    chunk_id = 0
    while start < len(words):
        end = min(start + CHUNK_SIZE, len(words))
        chunk_text = " ".join(words[start:end])
        chunk_hash = hashlib.md5(f"{source}_{chunk_id}".encode()).hexdigest()
        chunks.append(
            {
                "id": chunk_hash,
                "text": chunk_text,
                "source": source,
                "chunk_index": chunk_id,
            }
        )
        chunk_id += 1
        if end >= len(words):
            break
        start = end - CHUNK_OVERLAP

    return chunks


def _extract_text_from_file(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".txt":
            return path.read_text("utf-8", errors="ignore")
        elif ext == ".md":
            return path.read_text("utf-8", errors="ignore")
        elif ext == ".pdf":
            try:
                import PyPDF2

                text = ""
                with open(path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                return text
            except ImportError:
                return ""
        elif ext in (".docx", ".doc"):
            try:
                from docx import Document

                doc = Document(str(path))
                return "\n".join(p.text for p in doc.paragraphs)
            except ImportError:
                return ""
        elif ext == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(str(path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            except ImportError:
                return ""
        elif ext == ".csv":
            return path.read_text("utf-8", errors="ignore")
    except Exception as e:
        log.error("Error extracting %s: %s", path.name, e)
    return ""


def index_document(filepath: str) -> str:
    path = Path(filepath)
    if not path.exists():
        return f"Archivo no encontrado: {filepath}"

    text = _extract_text_from_file(path)
    if not text.strip():
        return f"No se pudo extraer texto de: {path.name}"

    chunks = _chunk_text(text, path.name)
    if not chunks:
        return f"Archivo vacío: {path.name}"

    collection = _get_collection()
    embedder = _get_embedder()

    existing = collection.get(ids=[c["id"] for c in chunks])
    if existing["ids"]:
        return f"'{path.name}' ya está indexado ({len(chunks)} chunks)"

    texts = [c["text"] for c in chunks]
    ids = [c["id"] for c in chunks]
    metadatas = [{"source": c["source"], "chunk_index": c["chunk_index"]} for c in chunks]

    embeddings = embedder.encode(texts, show_progress_bar=False).tolist()
    collection.add(ids=ids, embeddings=embeddings, documents=texts, metadatas=metadatas)

    return f"Indexado: {path.name} ({len(chunks)} chunks)"


def index_documents_folder() -> str:
    if not DOCUMENTS_DIR.exists():
        return "No hay carpeta de documentos."

    results = []
    for path in DOCUMENTS_DIR.iterdir():
        if path.is_file() and path.suffix.lower() in (".txt", ".md", ".pdf", ".docx", ".doc", ".pptx", ".csv"):
            result = index_document(str(path))
            results.append(result)

    return "\n".join(results) if results else "No se indexó ningún archivo."


def search_documents(query: str, top_k: int | None = None) -> str:
    if top_k is None:
        top_k = RAG_TOP_K

    collection = _get_collection()
    embedder = _get_embedder()

    count = collection.count()
    if count == 0:
        return "No hay documentos indexados. Usa 'indexa documentos' primero."

    query_embedding = embedder.encode([query], show_progress_bar=False).tolist()
    results = collection.query(
        query_embeddings=query_embedding,
        n_results=min(top_k, count),
    )

    if not results["documents"] or not results["documents"][0]:
        return "No encontré información relevante en los documentos."

    output = []
    for i, (doc, meta, score) in enumerate(zip(results["documents"][0], results["metadatas"][0], results["distances"][0])):
        source = meta.get("source", "desconocido")
        relevance = 1 - score
        output.append(f"[{source} (confianza: {relevance:.0%})] {doc[:300]}")

    return "\n\n".join(output)


def get_rag_context(query: str, top_k: int = 3) -> str:
    collection = _get_collection()
    embedder = _get_embedder()

    count = collection.count()
    if count == 0:
        return ""

    query_embedding = embedder.encode([query], show_progress_bar=False).tolist()
    results = collection.query(
        query_embeddings=query_embedding,
        n_results=min(top_k, count),
    )

    if not results["documents"] or not results["documents"][0]:
        return ""

    context_parts = []
    for doc, meta in zip(results["documents"][0], results["metadatas"][0]):
        source = meta.get("source", "desconocido")
        context_parts.append(f"(De: {source}) {doc}")

    return "\n\n".join(context_parts)
